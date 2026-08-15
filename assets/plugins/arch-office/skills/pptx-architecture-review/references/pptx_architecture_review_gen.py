#!/usr/bin/env python3
"""Скелет генератора презентации на архитектурный комитет (pptx, 16:9).

Запуск: python3 pptx_architecture_review_gen.py review.pptx
"""

import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

NAVY = RGBColor(0x1F, 0x38, 0x64)
GRAY = RGBColor(0x59, 0x59, 0x59)

# ── ЗАПОЛНИТЬ ────────────────────────────────────────────────────────────────
DECK_TITLE = ("Защита решения: <название>", "SAD v0.3 · Команда архитектуры · статус: на ревью")
# Слайды: (заголовок, [тезисы])
SLIDES = [
    ("Контекст и драйверы", [
        "Бизнес-цель → драйверы D1–D3",
        "Ограничения: стек, сроки, регуляторика",
    ]),
    ("Требования (NFR)", [
        "NFR-1 доступность 99.9% — проверка: синтетические пробы",
        "NFR-2 p95 ≤ 300 мс — проверка: нагрузочный тест",
        "NFR-3 RPO=0 для проводок — проверка: DR-учения",
    ]),
    ("Варианты и причины отказа", [
        "A — отклонён: не закрывает NFR-3 без дорогой доработки",
        "B — выбран: закрывает NFR-1..3 в штатном стеке",
        "0 — статус-кво: долг растёт, инцидент INC-2211 как прецедент",
    ]),
    ("Решение: C4-контекст", [
        "Схема: система ↔ смежники (3–5 блоков, см. SAD раздел 5)",
        "Границы ответственности контуров",
    ]),
    ("C4-контейнеры и ключевой поток", [
        "Контейнеры: api, worker, store; sequence потока create→confirm",
        "Инвариант: идемпотентность по tx_id на всех записях",
    ]),
    ("Стыки и контракты", [
        "IS-12 каналы→API; IS-13 шина событий — версии контрактов зафиксированы",
    ]),
    ("Что может сломаться", [
        "Риск: рост очереди при деградации СБП → митигация: backpressure + shed",
        "Открытая точка: лимиты хранилища событий — оценка к H1",
    ]),
    ("План внедрения и отката", [
        "Этапы с критериями выхода; откат каждого этапа — флагом маршрутизации",
    ]),
    ("Вопросы к комитету", [
        "Утвердить решение в объёме SAD v0.3",
        "Разрешить пилот на домене X в окне H1",
        "Открытая точка: лимиты хранилища — приемлем ли расчёт к H1?",
    ]),
]
# ─────────────────────────────────────────────────────────────────────────────

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


def add_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    band = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Emu(900000))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    p = band.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box = slide.shapes.add_textbox(Emu(600000), Emu(1200000), Emu(11000000), Emu(5200000))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▪ {bullet}"
        p.font.size = Pt(18)
        p.space_after = Pt(10)


def main(path: str) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(600000), Emu(2400000), Emu(11000000), Emu(2000000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = DECK_TITLE[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = DECK_TITLE[1]
    p2.font.size = Pt(16)
    p2.font.color.rgb = GRAY
    for title, bullets in SLIDES:
        add_slide(prs, title, bullets)
    prs.save(path)
    print(f"записан {path}: {len(prs.slides._sldIdLst)} слайдов")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "review.pptx")
