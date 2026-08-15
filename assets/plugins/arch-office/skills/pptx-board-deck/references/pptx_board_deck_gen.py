#!/usr/bin/env python3
"""Скелет генератора презентации для правления (pptx, 16:9).

Заполнение: DECK_TITLE и SLIDES (action title + тезисы + опциональная сноска).
Запуск: python3 pptx_board_deck_gen.py deck.pptx
"""

import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

NAVY = RGBColor(0x1F, 0x38, 0x64)
GRAY = RGBColor(0x59, 0x59, 0x59)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)

# ── ЗАПОЛНИТЬ ────────────────────────────────────────────────────────────────
DECK_TITLE = ("<Тема инициативы>", "Команда архитектуры · <дата> · статус: на утверждение")
# Слайды: (action title, [тезисы], сноска с допущениями или "")
SLIDES = [
    ("Проблема: <что болит и почему сейчас>", [
        "Цифра боли: простой/потери/риск — X в год",
        "Точка невозврата: почему откладывать дороже",
    ], ""),
    ("Предлагается: <суть решения одной строкой>", [
        "Что меняется для бизнеса — возможностями, не продуктами",
        "Что НЕ меняется (границы инициативы)",
    ], ""),
    ("Рассмотрены три варианта; рекомендуется B", [
        "A — почему отклонён (цена/срок/риск)",
        "B — рекомендуемый: суть в одну строку",
        "0 — «ничего не делать»: цена бездействия",
    ], ""),
    ("Как это работает", [
        "Схема 3–5 блоков: контур → поток → точка контроля",
        "Ключевой инвариант решения одной строкой",
    ], ""),
    ("Экономика: окупаемость 14 месяцев", [
        "Эффект: ₽X млн/год; затраты: ₽Y млн разово + ₽Z млн/год",
    ], "Допущения: ставка FTE, объём операций 2026 г., курсы на 01.08.2026"),
    ("Риски перемен и бездействия взвешены", [
        "Главный риск перемен — и его митигация",
        "Главный риск бездействия — и его цена",
    ], ""),
    ("Дорожная карта: три горизонта", [
        "H1 — пилот (результат, критерий перехода)",
        "H2 — промышленное распространение",
        "H3 — полный перевод и вывод legacy",
    ], ""),
    ("Запрашиваемое решение", [
        "Утвердить: целевое состояние и горизонт H1",
        "Профинансировать: ₽Y млн на H1",
        "Поручить: <кому> — отчёт о пилоте к <дата>",
    ], ""),
]
# ─────────────────────────────────────────────────────────────────────────────

SLIDE_W = Emu(12192000)  # 16:9
SLIDE_H = Emu(6858000)


def add_text_slide(prs: Presentation, title: str, bullets: list[str], footnote: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # пустой макет
    # Полоса заголовка.
    band = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Emu(900000))  # 1 = rectangle
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    tf = band.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Тезисы.
    box = slide.shapes.add_textbox(Emu(600000), Emu(1200000), Emu(11000000), Emu(4800000))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▪ {bullet}"
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    # Сноска с допущениями.
    if footnote:
        note = slide.shapes.add_textbox(Emu(600000), Emu(6100000), Emu(11000000), Emu(500000))
        p = note.text_frame.paragraphs[0]
        p.text = footnote
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY


def main(path: str) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Титул.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(600000), Emu(2400000), Emu(11000000), Emu(2000000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = DECK_TITLE[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = DECK_TITLE[1]
    p2.font.size = Pt(16)
    p2.font.color.rgb = GRAY

    for title, bullets, footnote in SLIDES:
        add_text_slide(prs, title, bullets, footnote)

    prs.save(path)
    print(f"записан {path}: {len(prs.slides.__iter__.__self__._sldIdLst)} слайдов")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "deck.pptx")
